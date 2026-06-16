"""
Team Blackbox 최종 발표자료 생성 스크립트
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── 색상 팔레트 (라이트 테마) ────────────────────────────
DARK_BG    = RGBColor(0xF1, 0xF5, 0xF9)   # 슬라이드 배경 (연한 회청색)
ACCENT     = RGBColor(0x05, 0x96, 0x69)   # 진한 민트 (라이트 배경용)
ACCENT2    = RGBColor(0x02, 0x84, 0xC7)   # 진한 하늘색
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x47, 0x5A, 0x6B)   # 본문 보조 텍스트
DARK_TEXT  = RGBColor(0x0F, 0x17, 0x2A)   # 기본 텍스트 (진남색)
CARD_BG    = RGBColor(0xFF, 0xFF, 0xFF)   # 카드 배경 (흰색)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 완전 빈 슬라이드

# ── 헬퍼 함수 ────────────────────────────────────────────

def add_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return sl

def txbox(sl, text, l, t, w, h, size=18, bold=False, color=DARK_TEXT,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def rect(sl, l, t, w, h, fill=CARD_BG, line=None):
    sh = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.color.rgb = RGBColor(0xE2,0xE8,0xF0)
        sh.line.width = Pt(0.5)
    return sh

def accent_bar(sl, t=1.05, color=ACCENT):
    sh = sl.shapes.add_shape(1, Inches(0.5), Inches(t), Inches(0.06), Inches(0.45))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()

def slide_title(sl, title, subtitle=None):
    accent_bar(sl)
    txbox(sl, title, 0.7, 0.9, 11, 0.6, size=28, bold=True, color=DARK_TEXT)
    if subtitle:
        txbox(sl, subtitle, 0.7, 1.5, 11, 0.4, size=14, color=LIGHT_GRAY)

def bullet_box(sl, items, l, t, w, h, size=14, color=WHITE, title=None, title_color=ACCENT):
    if title:
        txbox(sl, title, l, t, w, 0.35, size=13, bold=True, color=title_color)
        t += 0.38
        h -= 0.38
    card = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = card.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color

def chip(sl, text, l, t, w=1.5, h=0.32, bg=ACCENT, tc=DARK_TEXT):
    sh = sl.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))  # 9=rounded rect
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg
    sh.line.fill.background()
    tf = sh.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = tc

# ══════════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ══════════════════════════════════════════════════════════
sl = add_slide()

# 데코 원
for (cx, cy, size, alpha) in [(10.5, 1.5, 3.5, 0x1A6EE7B7), (1.0, 6.0, 2.0, 0x1A38BDF8)]:
    sh = sl.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(size), Inches(size))
    sh.fill.solid()
    sh.fill.fore_color.rgb = ACCENT
    sh.line.fill.background()

txbox(sl, "🔲  BLACKBOX", 1.0, 1.6, 11, 0.8, size=15, color=ACCENT, bold=True)
txbox(sl, "팀 프로젝트 기여도 추적\nEdTech SaaS 플랫폼", 1.0, 2.4, 11, 1.4,
      size=40, bold=True, color=DARK_TEXT)
txbox(sl, "컴퓨터소프트웨어과  |  프로젝트 구현  |  Team Blackbox",
      1.0, 4.1, 11, 0.5, size=15, color=LIGHT_GRAY)
txbox(sl, "2025", 1.0, 4.7, 4, 0.5, size=13, color=ACCENT2)

# ══════════════════════════════════════════════════════════
# 슬라이드 2 — 목차
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "목차  INDEX")

sections = [
    ("01", "프로젝트 개요", "팀 소개 · 프로젝트 소개 · 개발 일정"),
    ("02", "프로젝트 설계", "아키텍처 · 유스케이스 · 스토리보드\n화면 설계 · DB 설계"),
    ("03", "프로젝트 구현", "개발환경 · 구현 기능 ①②③\nAI 활용 기능"),
    ("04", "2학기 진행 계획", "기능 확장 계획 · 일정"),
    ("05", "결과 및 고찰", "소감 · 시행착오 3가지 · 극복 사례"),
]

for i, (num, title, desc) in enumerate(sections):
    col = i % 3
    row = i // 3
    lx = 0.4 + col * 4.3
    ty = 1.9 + row * 2.2
    rect(sl, lx, ty, 3.9, 1.9, fill=CARD_BG, line=ACCENT if i == 0 else None)
    txbox(sl, num, lx+0.15, ty+0.1, 0.8, 0.5, size=22, bold=True, color=ACCENT)
    txbox(sl, title, lx+0.15, ty+0.55, 3.5, 0.4, size=15, bold=True, color=DARK_TEXT)
    txbox(sl, desc, lx+0.15, ty+1.0, 3.6, 0.7, size=11, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════
# 슬라이드 3 — 팀 소개
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "01  프로젝트 개요", "팀 소개")

rect(sl, 0.4, 1.8, 12.5, 0.45, fill=RGBColor(0xE8,0xEF,0xF5))
txbox(sl, "팀원", 0.5, 1.85, 2, 0.35, size=12, bold=True, color=ACCENT)
txbox(sl, "역할", 4.0, 1.85, 4, 0.35, size=12, bold=True, color=ACCENT)
txbox(sl, "기여도", 10.5, 1.85, 2, 0.35, size=12, bold=True, color=ACCENT)

members = [
    ("홍길동 (팀장)", "백엔드 서버 구현, DB 설계, 깃허브·Google 연동", "35%"),
    ("팀원 B",        "프론트엔드(Next.js) 구현, UI/UX 설계",           "30%"),
    ("팀원 C",        "Notion 연동, 캘린더 기능, 파일 검사 기능",        "20%"),
    ("팀원 D",        "AI 분석 연동, Discord 알림, PDF 보고서",           "15%"),
]

for i, (name, role, pct) in enumerate(members):
    ty = 2.35 + i * 0.82
    bg = RGBColor(0xF0,0xF4,0xF8) if i % 2 == 0 else CARD_BG
    rect(sl, 0.4, ty, 12.5, 0.75, fill=bg)
    txbox(sl, name, 0.55, ty+0.18, 3.2, 0.5, size=13, bold=(i==0), color=DARK_TEXT)
    txbox(sl, role, 4.0, ty+0.18, 7.0, 0.5, size=12, color=LIGHT_GRAY)
    txbox(sl, pct,  10.5, ty+0.18, 2.0, 0.5, size=13, bold=True, color=ACCENT)

txbox(sl, "※ 팀원 이름 및 기여도는 실제 정보로 수정 필요",
      0.4, 6.9, 12, 0.4, size=10, color=LIGHT_GRAY, italic=True)

# ══════════════════════════════════════════════════════════
# 슬라이드 4 — 프로젝트 소개
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "01  프로젝트 개요", "프로젝트 소개")

rect(sl, 0.4, 1.8, 5.8, 1.1, fill=CARD_BG)
txbox(sl, "프로젝트명", 0.6, 1.85, 5, 0.35, size=11, color=ACCENT)
txbox(sl, "Blackbox", 0.6, 2.2, 5.5, 0.5, size=22, bold=True, color=DARK_TEXT)

rect(sl, 6.5, 1.8, 6.3, 1.1, fill=CARD_BG)
txbox(sl, "개발 기간", 6.7, 1.85, 5, 0.35, size=11, color=ACCENT)
txbox(sl, "2025년 3월 ~ 6월  (약 15주)", 6.7, 2.2, 6.0, 0.5, size=14, bold=True, color=DARK_TEXT)

rect(sl, 0.4, 3.1, 12.4, 1.35, fill=CARD_BG)
txbox(sl, "선정 이유", 0.6, 3.15, 11, 0.35, size=11, bold=True, color=ACCENT)
txbox(sl,
      "대학교 팀 프로젝트에서 '무임승차' 문제가 빈번하게 발생하지만, 기여도를 객관적으로 측정하는 도구가 없음.\n"
      "GitHub 커밋·Notion 편집·Google Docs 활동·회의 참석·태스크 완료 등 다양한 데이터를 자동 수집·정량화하여\n"
      "교수·팀원 모두가 납득할 수 있는 투명한 기여도 평가 시스템을 제공하고자 선정.",
      0.6, 3.55, 12.0, 0.85, size=12, color=DARK_TEXT)

features = [
    ("📊", "실시간\n기여도 측정", "Task·Meeting·File·\nGitHub·Notion·Google"),
    ("🤖", "AI 팀 분석", "GPT-4.1-nano 기반\n팀 역량 진단"),
    ("🔔", "스마트 경보", "불균형·과부하·\n마감임박 자동 감지"),
    ("📄", "PDF 보고서", "교수 제출용\n한글 지원 PDF"),
    ("🔗", "외부 연동", "GitHub·Notion·\nGoogle·Discord"),
]
for i, (icon, title, desc) in enumerate(features):
    lx = 0.4 + i * 2.48
    rect(sl, lx, 4.6, 2.3, 2.6, fill=RGBColor(0xF0,0xF4,0xF8), line=ACCENT2)
    txbox(sl, icon, lx+0.8, 4.75, 0.8, 0.5, size=22)
    txbox(sl, title, lx+0.1, 5.35, 2.1, 0.55, size=12, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    txbox(sl, desc,  lx+0.1, 5.95, 2.1, 0.9,  size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# 슬라이드 5 — 개발 일정
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "01  프로젝트 개요", "개발 일정")

weeks = [
    ("1~2주차",  "주제 선정 및 요구사항 분석",          "기획"),
    ("3~4주차",  "ERD 설계, DB 구축, API 명세 작성",    "설계"),
    ("5~7주차",  "백엔드 핵심 기능 구현 (인증·점수엔진·태스크)", "구현"),
    ("8~10주차", "프론트엔드 구현 (대시보드·보드·보고서)", "구현"),
    ("11~12주차","외부 연동 (GitHub·Notion·Google·Discord)", "연동"),
    ("13~14주차","AI 분석, 캘린더, 파일 검사 기능 추가",  "고도화"),
    ("15주차",   "최종 테스트 및 발표 준비",              "완료"),
]

PHASE_COLOR = {
    "기획":  RGBColor(0xF5,0x9E,0x0B),
    "설계":  RGBColor(0x38,0xBD,0xF8),
    "구현":  ACCENT,
    "연동":  RGBColor(0xA7,0x8B,0xFA),
    "고도화":RGBColor(0xFB,0x7B,0x5E),
    "완료":  RGBColor(0x34,0xD3,0x99),
}

for i, (week, task, phase) in enumerate(weeks):
    ty = 1.75 + i * 0.73
    bg = RGBColor(0xF0,0xF4,0xF8) if i % 2 == 0 else CARD_BG
    rect(sl, 0.4, ty, 12.5, 0.66, fill=bg)
    txbox(sl, week, 0.55, ty+0.14, 1.6, 0.38, size=12, bold=True, color=ACCENT2)
    txbox(sl, task, 2.3,  ty+0.14, 8.8, 0.38, size=12, color=DARK_TEXT)
    ph_color = PHASE_COLOR.get(phase, ACCENT)
    ph_sh = sl.shapes.add_shape(9, Inches(11.3), Inches(ty+0.16), Inches(1.4), Inches(0.32))
    ph_sh.fill.solid(); ph_sh.fill.fore_color.rgb = ph_color
    ph_sh.line.fill.background()
    tf = ph_sh.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = phase
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = DARK_TEXT

# ══════════════════════════════════════════════════════════
# 슬라이드 6 — 시스템 아키텍처
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "02  프로젝트 설계", "시스템 아키텍처")

def flow_box(sl, lx, ty, w, h, title, lines, bg, border=ACCENT):
    rect(sl, lx, ty, w, h, fill=bg, line=border)
    top = sl.shapes.add_shape(1, Inches(lx), Inches(ty), Inches(w), Inches(0.06))
    top.fill.solid(); top.fill.fore_color.rgb = border; top.line.fill.background()
    txbox(sl, title, lx+0.12, ty+0.1, w-0.2, 0.35, size=11, bold=True, color=DARK_TEXT)
    for i, line in enumerate(lines):
        txbox(sl, line, lx+0.12, ty+0.5+i*0.38, w-0.2, 0.35, size=10, color=LIGHT_GRAY)

def harrow(sl, x, y, label=""):
    """가로 흐름 화살표"""
    arr = sl.shapes.add_shape(1, Inches(x), Inches(y+0.1), Inches(0.35), Inches(0.06))
    arr.fill.solid(); arr.fill.fore_color.rgb = RGBColor(0x4A,0x5A,0x6A); arr.line.fill.background()
    tip = sl.shapes.add_shape(5, Inches(x+0.25), Inches(y+0.01), Inches(0.18), Inches(0.25))
    tip.fill.solid(); tip.fill.fore_color.rgb = RGBColor(0x4A,0x5A,0x6A); tip.line.fill.background()
    if label:
        txbox(sl, label, x-0.1, y+0.22, 0.7, 0.22, size=7, color=RGBColor(0x6A,0x7A,0x8A), align=PP_ALIGN.CENTER)

# 브라우저 (사용자)
rect(sl, 0.3, 2.1, 1.5, 1.9, fill=RGBColor(0xF0,0xF4,0xF8), line=RGBColor(0x38,0x4A,0x60))
txbox(sl, "👤\n사용자\n브라우저", 0.35, 2.25, 1.4, 1.6, size=11, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

harrow(sl, 1.82, 2.75, "HTTPS")

# Nginx
flow_box(sl, 2.2, 1.9, 1.5, 2.3, "🔀 Nginx",
    ["Reverse Proxy", "80 → 3000/8080", "정적 파일 서빙"],
    RGBColor(0xDC,0xFB,0xE8), ACCENT)

harrow(sl, 3.72, 2.75)

# Frontend
flow_box(sl, 4.1, 1.9, 2.3, 2.3, "🌐 Frontend",
    ["Next.js 14", "App Router", "TypeScript", "Tailwind CSS"],
    RGBColor(0xE0,0xF2,0xFE), ACCENT2)

harrow(sl, 6.42, 2.75, "REST API")

# Backend
flow_box(sl, 6.8, 1.9, 2.5, 2.3, "⚙️ Backend",
    ["Spring Boot 3.x", "Java 21", "Spring Security", "JWT Auth"],
    RGBColor(0xEC,0xFD,0xF5), ACCENT)

harrow(sl, 9.32, 2.75, "JDBC")

# DB
flow_box(sl, 9.7, 1.9, 2.0, 2.3, "🗄️ Database",
    ["PostgreSQL 15", "Flyway Migration", "21개 테이블"],
    RGBColor(0xF3,0xE8,0xFF), RGBColor(0xA7,0x8B,0xFA))

# Infra 하단
rect(sl, 0.3, 4.38, 11.4, 0.06, fill=RGBColor(0x05,0x96,0x69))
rect(sl, 0.3, 4.44, 11.4, 1.1, fill=RGBColor(0xF0,0xFD,0xF4), line=RGBColor(0x05,0x96,0x69))
txbox(sl, "🐳  Docker Compose  —  컨테이너 오케스트레이션  |  Volume 마운트  |  환경변수 관리",
      0.45, 4.52, 11.1, 0.35, size=10, bold=True, color=RGBColor(0x05,0x96,0x69))
txbox(sl, "blackbox-nginx    blackbox-frontend    blackbox-backend    blackbox-db    (공유 네트워크: blackbox-net)",
      0.45, 4.9, 11.1, 0.55, size=9, color=RGBColor(0x0F,0x17,0x2A))

# 외부 서비스 (우측 수직)
ext_list = [
    ("🐙 GitHub", RGBColor(0xE8,0xEB,0xEF), "OAuth·Webhook·REST"),
    ("📝 Notion", RGBColor(0xF5,0xF5,0xF3), "Integration Token"),
    ("📊 Google", RGBColor(0xE8,0xEF,0xF5), "OAuth·Drive·Docs"),
    ("💬 Discord", RGBColor(0xE9,0xE9,0xFF), "Webhook"),
    ("🤖 OpenAI", RGBColor(0xDC,0xFB,0xE8), "GPT-4.1-nano"),
]
txbox(sl, "외부 서비스", 11.55, 1.75, 1.7, 0.3, size=10, bold=True, color=ACCENT2)
for i, (name, bg, sub) in enumerate(ext_list):
    ty = 2.1 + i * 0.96
    rect(sl, 11.55, ty, 1.65, 0.82, fill=bg, line=RGBColor(0x3A,0x4A,0x5A))
    txbox(sl, name, 11.62, ty+0.06, 1.5, 0.32, size=9, bold=True, color=DARK_TEXT)
    txbox(sl, sub,  11.62, ty+0.42, 1.5, 0.28, size=7.5, color=LIGHT_GRAY)
    # 연결 점선
    arr2 = sl.shapes.add_shape(1, Inches(9.7+2.0), Inches(ty+0.35), Inches(1.55), Inches(0.03))
    arr2.fill.solid(); arr2.fill.fore_color.rgb = RGBColor(0x3A,0x4A,0x5A); arr2.line.fill.background()

# ══════════════════════════════════════════════════════════
# 슬라이드 7 — DB 설계
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "02  프로젝트 설계", "데이터베이스 설계")

tables = [
    ("users",               "사용자 (Google OAuth)", "id, name, email, role, picture"),
    ("projects",            "프로젝트",               "id, name, description, is_archived"),
    ("project_members",     "프로젝트 멤버",           "project_id, user_id, role"),
    ("activity_logs",       "활동 로그 (기여도 원천)", "id, project_id, user_id, event_type, score_weight, payload(JSONB)"),
    ("tasks",               "태스크",                  "id, project_id, title, status, priority, due_date"),
    ("meetings",            "회의록",                  "id, project_id, title, meeting_at, notes"),
    ("file_vaults",         "파일 검사",               "id, project_id, file_name, sha256_hash, version"),
    ("github_installations","GitHub 연동",             "id, project_id, repo_full_name, token"),
    ("github_user_mappings","GitHub 계정 매핑",        "id, project_id, user_id, github_username"),
    ("alerts",              "경보",                    "id, project_id, alert_type, severity, is_resolved"),
    ("project_alert_config","경보 설정",               "id, project_id, deadline_days, discord_webhook_url"),
    ("ai_analysis_results", "AI 분석 캐시",            "id, project_id, analysis_type, result, created_at"),
]

txbox(sl, f"총 21개 테이블  (Flyway V1~V21 마이그레이션)", 0.4, 1.75, 12, 0.4, size=12, color=ACCENT2)

cols = [[], [], []]
for i, t in enumerate(tables):
    cols[i % 3].append(t)

for ci, col in enumerate(cols):
    lx = 0.4 + ci * 4.3
    for ri, (name, desc, fields) in enumerate(col):
        ty = 2.2 + ri * 1.3
        rect(sl, lx, ty, 4.0, 1.2, fill=CARD_BG, line=ACCENT if ri==0 and ci==0 else RGBColor(0x2A,0x3A,0x4A))
        txbox(sl, name,   lx+0.15, ty+0.05, 3.7, 0.3,  size=11, bold=True, color=ACCENT)
        txbox(sl, desc,   lx+0.15, ty+0.35, 3.7, 0.25, size=10, color=DARK_TEXT)
        txbox(sl, fields, lx+0.15, ty+0.62, 3.7, 0.5,  size=9,  color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════
# 슬라이드 8 — 스토리보드
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "02  프로젝트 설계", "스토리보드 — 주요 화면 이동 경로")

def sb_box(sl, lx, ty, num, title, desc, color=ACCENT2):
    """스토리보드 화면 박스"""
    # 외곽
    outer = sl.shapes.add_shape(1, Inches(lx), Inches(ty), Inches(2.6), Inches(1.95))
    outer.fill.solid(); outer.fill.fore_color.rgb = CARD_BG
    outer.line.color.rgb = color; outer.line.width = Pt(1.5)
    # 상단 컬러 바
    top = sl.shapes.add_shape(1, Inches(lx), Inches(ty), Inches(2.6), Inches(0.35))
    top.fill.solid(); top.fill.fore_color.rgb = color; top.line.fill.background()
    # 번호 뱃지
    badge = sl.shapes.add_shape(9, Inches(lx+0.08), Inches(ty+0.06), Inches(0.26), Inches(0.22))
    badge.fill.solid(); badge.fill.fore_color.rgb = RGBColor(0x0F,0x17,0x2A); badge.line.fill.background()
    bt = badge.text_frame; bt.paragraphs[0].alignment = PP_ALIGN.CENTER
    br = bt.paragraphs[0].add_run(); br.text = num
    br.font.size = Pt(8); br.font.bold = True; br.font.color.rgb = WHITE
    # 제목
    txbox(sl, title, lx+0.38, ty+0.06, 2.1, 0.24, size=10, bold=True, color=DARK_TEXT)
    # 화면 내용 영역 (회색 박스)
    inner = sl.shapes.add_shape(1, Inches(lx+0.12), Inches(ty+0.43), Inches(2.36), Inches(1.05))
    inner.fill.solid(); inner.fill.fore_color.rgb = RGBColor(0xF0,0xF4,0xF8); inner.line.fill.background()
    txbox(sl, desc, lx+0.18, ty+0.48, 2.25, 0.95, size=9, color=LIGHT_GRAY)
    # 하단 URL 표시
    txbox(sl, f"/{''.join(title.split()[0]).lower()}", lx+0.12, ty+1.55, 2.3, 0.28,
          size=8, color=RGBColor(0x4A,0x5A,0x6A), italic=True)

def arr(sl, lx, ty, direction="→"):
    txbox(sl, direction, lx, ty, 0.32, 0.35, size=16, bold=True, color=RGBColor(0x4A,0x5A,0x6A), align=PP_ALIGN.CENTER)

# ── 행 1: 로그인 → 프로젝트 목록 → 프로젝트 홈 → 기여도 ──
R1 = 1.72
screens_r1 = [
    ("①", "로그인",       "Google 소셜 로그인\nJWT 토큰 발급",          ACCENT),
    ("②", "프로젝트 목록", "내 프로젝트 카드\n새 프로젝트 생성 버튼",    ACCENT2),
    ("③", "프로젝트 홈",  "팀원 목록·역할\n아카이브 / 설정 이동",       RGBColor(0xA7,0x8B,0xFA)),
    ("④", "기여도 대시보드","점수 차트·경보 현황\nPDF 보고서 다운로드",  RGBColor(0x34,0xD3,0x99)),
]
for i, (num, title, desc, color) in enumerate(screens_r1):
    lx = 0.35 + i * 3.12
    sb_box(sl, lx, R1, num, title, desc, color)
    if i < 3:
        arr(sl, lx + 2.63, R1 + 0.78)

# 행1 → 행2 연결선 (④ 아래로)
arr(sl, 0.35 + 3*3.12 + 1.14, R1 + 1.98, "↓")

# ── 행 2: 태스크 보드 ← 캘린더 ← AI분석 ← 설정 ──
R2 = 4.3
screens_r2 = [
    ("⑤", "태스크 보드",  "Kanban 보드\n우선순위별 색 구분",            RGBColor(0xFB,0x7B,0x5E)),
    ("⑥", "캘린더",       "월별 그리드 뷰\n태스크·회의 일정 표시",      RGBColor(0x38,0xBD,0xF8)),
    ("⑦", "AI 분석",      "팀 역량 진단\n커밋 품질 점수화",              RGBColor(0xF5,0x9E,0x0B)),
    ("⑧", "설정",         "GitHub·Notion·Google\nDiscord Webhook 연동",  RGBColor(0x94,0xA3,0xB8)),
]
for i, (num, title, desc, color) in enumerate(screens_r2):
    lx = 0.35 + i * 3.12
    sb_box(sl, lx, R2, num, title, desc, color)
    if i < 3:
        arr(sl, lx + 2.63, R2 + 0.78)

# 행2 화살표 방향 표기 (← 순서)
txbox(sl, "※ ②~④ 는 상단 네비게이션 탭으로 언제든 이동 가능",
      0.35, 6.55, 12.5, 0.35, size=10, color=LIGHT_GRAY, italic=True)

# ══════════════════════════════════════════════════════════
# 슬라이드 8-1 — UseCase Diagram
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "02  프로젝트 설계", "유스케이스 다이어그램")

def uc_oval(sl, lx, ty, w, h, label, color):
    uc = sl.shapes.add_shape(9, Inches(lx), Inches(ty), Inches(w), Inches(h))
    uc.fill.solid(); uc.fill.fore_color.rgb = RGBColor(0xF0,0xF8,0xFF)
    uc.line.color.rgb = color; uc.line.width = Pt(1.5)
    tf = uc.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = label
    r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = color

def actor_card(sl, lx, ty, icon, label, color):
    sh = sl.shapes.add_shape(9, Inches(lx), Inches(ty), Inches(1.4), Inches(1.0))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xF0,0xF8,0xFF)
    sh.line.color.rgb = color; sh.line.width = Pt(2)
    txbox(sl, icon,  lx+0.48, ty+0.05, 0.5, 0.38, size=18, color=color, align=PP_ALIGN.CENTER)
    txbox(sl, label, lx,      ty+0.6,  1.4, 0.35, size=9, bold=True, color=color, align=PP_ALIGN.CENTER)

# 시스템 경계
sb = sl.shapes.add_shape(1, Inches(1.7), Inches(1.72), Inches(9.55), Inches(5.5))
sb.fill.solid(); sb.fill.fore_color.rgb = RGBColor(0xF8,0xFA,0xFC)
sb.line.color.rgb = ACCENT2; sb.line.width = Pt(2)
txbox(sl, "《 Blackbox 시스템 》", 1.75, 1.74, 4, 0.28, size=10, bold=True, color=ACCENT2)

# ── 그룹 배경 ──
groups = [
    (1.82, 1.98, 2.9, 5.1, RGBColor(0xEE,0xFB,0xF5), ACCENT,                    "인증 · 프로젝트"),
    (4.85, 1.98, 2.9, 5.1, RGBColor(0xE0,0xF2,0xFE), ACCENT2,                   "활동 관리"),
    (7.88, 1.98, 3.2, 5.1, RGBColor(0xF3,0xE8,0xFF), RGBColor(0xA7,0x8B,0xFA),  "분석 · 보고서"),
]
for lx, ty, w, h, bg, border, glabel in groups:
    rect(sl, lx, ty, w, h, fill=bg)
    txbox(sl, glabel, lx+0.1, ty+0.04, w-0.15, 0.25,
          size=8, bold=True, color=border, italic=True)

# ── 유스케이스 ──
C1 = ACCENT; C2 = ACCENT2; C3 = RGBColor(0xA7,0x8B,0xFA)
C4 = RGBColor(0xF5,0x9E,0x0B); C5 = RGBColor(0xEF,0x44,0x44); C6 = RGBColor(0x34,0xD3,0x99)

col1_ucs = ["Google 로그인", "프로젝트 생성", "프로젝트 참여", "멤버 역할 변경", "프로젝트 아카이브"]
col1_clr = [C1, C1, C1, C4, C4]
col2_ucs = ["태스크 관리", "회의록 작성", "파일 업로드·검증", "캘린더 일정 등록", "기여도 점수 조회"]
col3_ucs = ["AI 팀 분석", "커밋 품질 분석", "경보 확인·설정", "PDF 보고서 출력", "외부 서비스 연동"]
col3_clr = [C3, C3, C5, C6, C6]

for i, (label, color) in enumerate(zip(col1_ucs, col1_clr)):
    uc_oval(sl, 1.88, 2.28+i*0.9, 2.75, 0.72, label, color)
for i, label in enumerate(col2_ucs):
    uc_oval(sl, 4.92, 2.28+i*0.9, 2.75, 0.72, label, C2)
for i, (label, color) in enumerate(zip(col3_ucs, col3_clr)):
    uc_oval(sl, 7.95, 2.28+i*0.9, 3.05, 0.72, label, color)

# ── 액터 카드 ──
actor_card(sl, 0.15, 2.35, "👤", "팀원 (학생)", ACCENT2)
actor_card(sl, 0.15, 4.15, "👑", "팀장",         C4)
actor_card(sl, 11.4, 2.35, "🎓", "교수 / TA",    C3)
actor_card(sl, 11.4, 4.15, "🔗", "외부 시스템",  C6)

# ── 연결 표시 (단순 수평선) ──
def hline(sl, x1, x2, y, color):
    ln = sl.shapes.add_shape(1, Inches(x1), Inches(y), Inches(x2-x1), Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = color; ln.line.fill.background()

# 팀원 → 열1, 열2 (수평선으로 연결)
hline(sl, 1.55, 1.88, 2.64, ACCENT2)   # → 로그인
hline(sl, 1.55, 1.88, 3.54, ACCENT2)   # → 생성
hline(sl, 1.55, 1.88, 4.44, ACCENT2)   # → 참여
hline(sl, 1.55, 4.92, 5.34, ACCENT2)   # → 태스크
hline(sl, 1.55, 4.92, 6.24, ACCENT2)   # → 회의록
# 팀장 → 팀장 전용
hline(sl, 1.55, 1.88, 4.59, C4)        # → 멤버역할변경
hline(sl, 1.55, 1.88, 5.49, C4)        # → 아카이브
# 교수 → 분석
hline(sl, 11.0, 11.4, 2.64, C3)
hline(sl, 11.0, 11.4, 3.54, C3)
hline(sl, 11.0, 11.4, 4.44, C5)
hline(sl, 11.0, 11.4, 5.34, C6)

# ── 범례 (하단) ──
legend_items2 = [
    (C1, "인증·프로젝트"), (C2, "활동 관리"), (C3, "AI 분석"),
    (C4, "팀장 전용"),     (C5, "경보"),      (C6, "보고서·연동"),
]
for i, (color, label) in enumerate(legend_items2):
    lx = 1.7 + i * 1.9
    dot = sl.shapes.add_shape(9, Inches(lx), Inches(7.12), Inches(0.18), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
    txbox(sl, label, lx+0.22, 7.1, 1.65, 0.25, size=9, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════
# 슬라이드 8-2 — 핵심 화면 설계
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "02  프로젝트 설계", "화면 설계 — 핵심 화면 소개")

# 기여도 대시보드 큰 와이어프레임
rect(sl, 0.4, 1.78, 6.0, 5.0, fill=RGBColor(0xF8,0xFA,0xFC), line=RGBColor(0x38,0x4A,0x60))
bar2 = sl.shapes.add_shape(1, Inches(0.4), Inches(1.78), Inches(6.0), Inches(0.3))
bar2.fill.solid(); bar2.fill.fore_color.rgb = RGBColor(0xBF,0xDB,0xF2); bar2.line.fill.background()
txbox(sl, "기여도 대시보드", 0.55, 1.82, 4, 0.22, size=8, bold=True, color=ACCENT2)

# 점수 카드들
score_cards = [("홍길동", "87", ACCENT), ("팀원B", "72", ACCENT2), ("팀원C", "51", RGBColor(0xF5,0x9E,0x0B))]
for i, (name, score, color) in enumerate(score_cards):
    cx = 0.55 + i * 1.9
    rect(sl, cx, 2.15, 1.7, 1.0, fill=RGBColor(0xFF,0xFF,0xFF), line=color)
    txbox(sl, name,  cx+0.1, 2.2,  1.5, 0.3, size=9, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    txbox(sl, score, cx+0.1, 2.52, 1.5, 0.45, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
    txbox(sl, "점",  cx+0.1, 2.98, 1.5, 0.18, size=8, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 막대 그래프 영역
rect(sl, 0.55, 3.3, 5.6, 1.6, fill=RGBColor(0xF0,0xF4,0xF8))
txbox(sl, "활동 유형별 기여도", 0.65, 3.35, 4, 0.25, size=8, color=LIGHT_GRAY)
bar_labels = ["Task", "Meet", "File", "GitHub", "Notion", "Google"]
bar_vals   = [0.85, 0.65, 0.45, 0.78, 0.52, 0.38]
for i, (lbl, val) in enumerate(zip(bar_labels, bar_vals)):
    bx = 0.65 + i * 0.9
    bh = val * 1.0
    rect(sl, bx, 3.6+(1.0-bh), 0.65, bh, fill=ACCENT if val > 0.6 else RGBColor(0x38,0xBD,0xF8))
    txbox(sl, lbl, bx-0.05, 4.68, 0.75, 0.2, size=7, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 경보 영역
rect(sl, 0.55, 5.05, 5.6, 0.55, fill=RGBColor(0xFE,0xF2,0xF2), line=RGBColor(0xEF,0x44,0x44))
txbox(sl, "⚠  IMBALANCE 경보: 팀원 간 기여도 편차 과도 (기준 30pt 초과)",
      0.7, 5.1, 5.3, 0.38, size=9, color=RGBColor(0xB9,0x1C,0x1C))

txbox(sl, "〈 기여도 대시보드 〉", 0.4, 6.7, 6.0, 0.38, size=10, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# 태스크 보드 와이어프레임
rect(sl, 6.8, 1.78, 6.1, 5.0, fill=RGBColor(0xF8,0xFA,0xFC), line=RGBColor(0x38,0x4A,0x60))
bar3 = sl.shapes.add_shape(1, Inches(6.8), Inches(1.78), Inches(6.1), Inches(0.3))
bar3.fill.solid(); bar3.fill.fore_color.rgb = RGBColor(0xBF,0xDB,0xF2); bar3.line.fill.background()
txbox(sl, "태스크 보드 (Kanban)", 6.95, 1.82, 4, 0.22, size=8, bold=True, color=ACCENT2)

# 필터 버튼들
filter_btns = [("전체", ACCENT), ("긴급", RGBColor(0xD9,0x46,0xEF)),
               ("높음", RGBColor(0xF9,0x73,0x16)), ("보통", RGBColor(0xEA,0xB3,0x08))]
for i, (lbl, col) in enumerate(filter_btns):
    fb = sl.shapes.add_shape(9, Inches(6.95+i*1.15), Inches(2.15), Inches(1.0), Inches(0.26))
    fb.fill.solid(); fb.fill.fore_color.rgb = col if i==0 else RGBColor(0xFF,0xFF,0xFF)
    fb.line.fill.background()
    ft = fb.text_frame; ft.paragraphs[0].alignment = PP_ALIGN.CENTER
    fr = ft.paragraphs[0].add_run(); fr.text = lbl
    fr.font.size = Pt(8); fr.font.color.rgb = DARK_TEXT if i==0 else col

# 칸반 컬럼
kanban_cols = [("TODO", [("🔴 로그인 구현", "긴급"), ("🟠 API 설계", "높음")]),
               ("진행중", [("🟡 DB 설계", "보통")]),
               ("완료", [("✅ 기획서 작성", "낮음"), ("✅ ERD 완성", "낮음")])]
for ci, (col_title, cards) in enumerate(kanban_cols):
    cx = 6.95 + ci * 1.95
    rect(sl, cx, 2.5, 1.8, 0.28, fill=RGBColor(0xE2,0xE8,0xF0))
    txbox(sl, col_title, cx+0.05, 2.54, 1.7, 0.2, size=8, bold=True, color=DARK_TEXT)
    for ri, (card_title, priority) in enumerate(cards):
        rect(sl, cx+0.05, 2.86+ri*1.05, 1.7, 0.92, fill=RGBColor(0xF0,0xF4,0xF8), line=RGBColor(0x2A,0x3A,0x4A))
        txbox(sl, card_title, cx+0.12, 2.9+ri*1.05, 1.55, 0.4, size=7.5, color=DARK_TEXT)
        txbox(sl, priority,   cx+0.12, 3.3+ri*1.05, 1.0, 0.2, size=7, color=LIGHT_GRAY)

txbox(sl, "〈 태스크 보드 (Kanban) 〉", 6.8, 6.7, 6.1, 0.38, size=10, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# 슬라이드 8-3 — 화면 설계 ② (회의록·AI분석·파일검사·설정)
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "02  프로젝트 설계", "화면 설계 — 주요 화면 소개 ②")

def mini_screen(sl, lx, ty, w, h, title, rows, label_bottom):
    """미니 화면 와이어프레임"""
    outer = sl.shapes.add_shape(1, Inches(lx), Inches(ty), Inches(w), Inches(h))
    outer.fill.solid(); outer.fill.fore_color.rgb = RGBColor(0xF8,0xFA,0xFC)
    outer.line.color.rgb = RGBColor(0x38,0x4A,0x60); outer.line.width = Pt(1)
    # 타이틀바
    bar = sl.shapes.add_shape(1, Inches(lx), Inches(ty), Inches(w), Inches(0.27))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0xBF,0xDB,0xF2); bar.line.fill.background()
    txbox(sl, title, lx+0.1, ty+0.04, w-0.15, 0.2, size=8, bold=True, color=ACCENT2)
    # 내용 행
    row_h = (h - 0.32) / max(len(rows), 1)
    for i, (row_txt, row_color) in enumerate(rows):
        ry = ty + 0.32 + i * row_h
        if row_color:
            rb = sl.shapes.add_shape(1, Inches(lx+0.08), Inches(ry+0.03),
                                     Inches(w-0.16), Inches(row_h-0.06))
            rb.fill.solid(); rb.fill.fore_color.rgb = row_color; rb.line.fill.background()
        txbox(sl, row_txt, lx+0.12, ry+0.05, w-0.22, row_h-0.1, size=8, color=DARK_TEXT)
    txbox(sl, label_bottom, lx, ty+h+0.05, w, 0.28,
          size=10, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# ── 4개 화면을 2x2 그리드로 ────────────────────────────────
SW = 6.1; SH = 2.4   # 각 화면 크기

# 1) 회의록
mini_screen(sl, 0.35, 1.72, SW, SH, "📝  회의록 작성",
    [("제목   2025-06-01 스프린트 회의",          RGBColor(0xFF,0xFF,0xFF)),
     ("일시   2025-06-01  |  장소  Discord",       RGBColor(0xF0,0xF4,0xF8)),
     ("목적   기여도 알고리즘 최종 확정",           RGBColor(0xF0,0xF4,0xF8)),
     ("참석자  홍길동 ✓  팀원B ✓  팀원C ✓",        RGBColor(0xE8,0xEF,0xF5)),
     ("[✉ Google Docs 내보내기]     [저장]",        RGBColor(0xD1,0xFA,0xE5))],
    "〈 회의록 작성 화면 〉")

# 2) AI 분석
mini_screen(sl, 6.85, 1.72, SW, SH, "🤖  AI 팀 분석",
    [("[팀 역량 분석]        [커밋 품질 분석]",    RGBColor(0xFF,0xFF,0xFF)),
     ("✅ 강점  회의 참석률 95% (매우 우수)",        RGBColor(0xDC,0xFB,0xE8)),
     ("⚠️ 약점  GitHub 커밋 품질 편차 큼",          RGBColor(0xFE,0xE2,0xE2)),
     ("💡 제안  코드 리뷰 문화 도입 권장",           RGBColor(0xF0,0xF4,0xF8)),
     ("커밋 품질 평균  ★★★★☆  3.8 / 5.0",         RGBColor(0xF0,0xFD,0xF4))],
    "〈 AI 분석 화면 〉")

# 3) 파일 검사
mini_screen(sl, 0.35, 4.42, SW, SH, "🔒  파일 검사 (File Vault)",
    [("파일명             버전   SHA-256      상태",  RGBColor(0xE8,0xEF,0xF5)),
     ("report_final.pdf   v3     a1b2c3…      ✅ 정상", RGBColor(0xFF,0xFF,0xFF)),
     ("design_v2.fig      v1     ff00aa…      ✅ 정상", RGBColor(0xF0,0xF4,0xF8)),
     ("old_backup.zip     v2     dead01…      ⚠️ 변조", RGBColor(0xFE,0xE2,0xE2)),
     ("[파일 업로드]   [무결성 검증]   [Drive 전송]", RGBColor(0xD1,0xFA,0xE5))],
    "〈 파일 검사 화면 〉")

# 4) 설정 (외부 연동)
mini_screen(sl, 6.85, 4.42, SW, SH, "⚙️  설정 — 외부 연동",
    [("[GitHub]  [Notion]  [Google]  [Discord]",      RGBColor(0xFF,0xFF,0xFF)),
     ("🐙 GitHub  repo: thswjdgy/blackbox   ✅ 연결됨", RGBColor(0xDC,0xFB,0xE8)),
     ("📝 Notion  Database 연동             ✅ 연결됨", RGBColor(0xDC,0xFB,0xE8)),
     ("📊 Google  Drive 폴더 연동           ✅ 연결됨", RGBColor(0xDC,0xFB,0xE8)),
     ("💬 Discord  Webhook URL  [저장] [테스트]",      RGBColor(0xF0,0xF4,0xF8))],
    "〈 설정 — 외부 연동 화면 〉")

# ══════════════════════════════════════════════════════════
# 슬라이드 9 — 개발환경 및 도구
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "03  프로젝트 구현", "개발환경 및 도구")

env_data = [
    ("구분", "기술 스택", "버전"),
    ("프론트엔드", "Next.js · TypeScript · Tailwind CSS", "14.x · 5.x · 3.4"),
    ("백엔드", "Spring Boot · Java · Spring Security", "3.x · 21 · 6.x"),
    ("데이터베이스", "PostgreSQL · Flyway Migration", "15 · 9.x"),
    ("인프라", "Docker Compose · Nginx", "Latest"),
    ("AI", "OpenAI GPT-4.1-nano (WebClient)", "API"),
    ("외부 연동", "GitHub OAuth App · Notion API · Google Drive API · Discord Webhook", "-"),
]

tool_data = [
    ("IDE", "IntelliJ IDEA · VS Code"),
    ("형상관리", "Git · GitHub"),
    ("빌드", "Gradle (백엔드) · npm (프론트엔드)"),
    ("DB 관리", "DBeaver · pgAdmin"),
    ("API 테스트", "Postman · Swagger UI"),
    ("협업", "GitHub Issues · Discord"),
]

# 개발환경 테이블
rect(sl, 0.4, 1.8, 7.7, 0.4, fill=RGBColor(0xEC,0xFD,0xF5))
for xi, col in enumerate(["구분", "기술 스택", "버전"]):
    lx = [0.5, 2.1, 7.0][xi]
    txbox(sl, col, lx, 1.85, 2.5, 0.3, size=11, bold=True, color=ACCENT)

for i, (div, stack, ver) in enumerate(env_data[1:]):
    ty = 2.25 + i * 0.68
    bg = RGBColor(0xF0,0xF4,0xF8) if i%2==0 else CARD_BG
    rect(sl, 0.4, ty, 7.7, 0.62, fill=bg)
    txbox(sl, div,   0.5,  ty+0.13, 1.5, 0.38, size=11, bold=True, color=DARK_TEXT)
    txbox(sl, stack, 2.1,  ty+0.13, 4.7, 0.38, size=11, color=LIGHT_GRAY)
    txbox(sl, ver,   7.0,  ty+0.13, 1.0, 0.38, size=10, color=ACCENT2)

# 개발도구
txbox(sl, "개발 도구", 8.4, 1.8, 4.7, 0.38, size=13, bold=True, color=ACCENT)
for i, (category, tools) in enumerate(tool_data):
    ty = 2.25 + i * 0.76
    rect(sl, 8.4, ty, 4.7, 0.68, fill=CARD_BG)
    txbox(sl, category, 8.55, ty+0.06, 1.5, 0.3, size=10, bold=True, color=ACCENT2)
    txbox(sl, tools,    8.55, ty+0.35, 4.5, 0.28, size=10, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════
# 슬라이드 10 — 구현 기능 (1) 기여도 측정
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "03  프로젝트 구현", "구현 기능 ① — 기여도 점수 측정 시스템")

rect(sl, 0.4, 1.8, 12.5, 1.5, fill=CARD_BG, line=ACCENT)
txbox(sl, "핵심 로직 — ScoreEngine", 0.6, 1.85, 6, 0.35, size=12, bold=True, color=ACCENT)
txbox(sl,
      "각 팀원의 7가지 활동 데이터를 수집 → raw 점수 산출 → Min-Max 정규화 → 0~100점 환산\n"
      "• 가중치: Task(30%) · Meeting(20%) · File(10%) · GitHub(20%) · Notion(10%) · Google(10%)\n"
      "• AI 커밋 품질 분석 결과(1~5점)를 scoreWeight에 반영 (x0.5 ~ x1.5)",
      0.6, 2.2, 12.1, 1.0, size=12, color=DARK_TEXT)

sources = [
    ("📋 Task",    "완료 태스크 수\n× 우선순위 가중치"),
    ("📅 Meeting", "회의 참석 횟수\n& 시간"),
    ("📁 File",    "파일 업로드 수\n& 검증 통과"),
    ("🐙 GitHub",  "커밋 수 + PR\n× 품질 점수"),
    ("📝 Notion",  "페이지 생성\n& 편집 횟수"),
    ("📊 Google",  "Docs/Sheets\n활동 횟수"),
]
for i, (icon_title, desc) in enumerate(sources):
    lx = 0.4 + i * 2.08
    rect(sl, lx, 3.45, 1.95, 1.7, fill=RGBColor(0xF0,0xF4,0xF8), line=ACCENT2)
    txbox(sl, icon_title, lx+0.1, 3.55, 1.75, 0.45, size=12, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    txbox(sl, desc,       lx+0.1, 4.05, 1.75, 0.9,  size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txbox(sl, "→", 12.6, 4.0, 0.5, 0.5, size=18, color=ACCENT2)

rect(sl, 0.4, 5.3, 12.5, 1.9, fill=RGBColor(0xEC,0xFD,0xF5))
txbox(sl, "핵심 코드 — getReport() 점수 산출", 0.6, 5.35, 8, 0.35, size=12, bold=True, color=ACCENT)
txbox(sl,
      "double[] rawScores = {taskScore, meetingScore, fileScore, githubScore, notionScore, googleScore};\n"
      "double[] weights   = {0.30, 0.20, 0.10, 0.20, 0.10, 0.10};\n"
      "// Min-Max 정규화 후 0~100 점 환산 → 최종 normalized score 산출",
      0.6, 5.75, 12.1, 1.3, size=10, color=RGBColor(0x86,0xEF,0xAC))

# ══════════════════════════════════════════════════════════
# 슬라이드 11 — 구현 기능 (2) 외부 연동
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "03  프로젝트 구현", "구현 기능 ② — 외부 서비스 연동")

integrations = [
    ("🐙  GitHub",
     ["OAuth App으로 GitHub 계정 자동 매핑",
      "Webhook 수신 → 커밋/PR 즉시 반영",
      "6시간 폴링 백업 (Webhook 누락 방지)",
      "repo URL 정규화 (https://github.com/ 자동 제거)"],
     ACCENT),
    ("📝  Notion",
     ["Notion Integration Token 등록",
      "Database 페이지 생성/편집 자동 감지",
      "6시간 폴링으로 활동 수집 → activity_log 저장",
      "Discord에 Notion 업데이트 알림 전송"],
     RGBColor(0xA7,0x8B,0xFA)),
    ("📊  Google",
     ["OAuth 2.0으로 Google 계정 연동",
      "Drive/Docs/Sheets 편집 활동 수집",
      "파일 업로드 시 Drive 자동 push",
      "회의록 → Google Docs 내보내기"],
     RGBColor(0xFB,0xBC,0x05)),
    ("💬  Discord",
     ["프로젝트별 Webhook URL 등록",
      "경보 발생 시 Embed 메시지 전송",
      "GitHub·Notion·Google 업데이트 알림",
      "태스크·회의·파일 생성 시 알림"],
     RGBColor(0x58,0x65,0xF2)),
]

for i, (title, bullets, color) in enumerate(integrations):
    col = i % 2; row = i // 2
    lx = 0.4 + col * 6.35
    ty = 1.85 + row * 2.65
    rect(sl, lx, ty, 6.1, 2.45, fill=CARD_BG)
    # 색 바
    sh = sl.shapes.add_shape(1, Inches(lx), Inches(ty), Inches(0.08), Inches(2.45))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    txbox(sl, title, lx+0.2, ty+0.1, 5.7, 0.45, size=15, bold=True, color=DARK_TEXT)
    for j, b in enumerate(bullets):
        txbox(sl, f"• {b}", lx+0.2, ty+0.6+j*0.44, 5.7, 0.4, size=11, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════
# 슬라이드 12 — 구현 기능 (3) 기타
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "03  프로젝트 구현", "구현 기능 ③ — 파일 검사 · 캘린더 · 경보 시스템")

feats = [
    ("🔒  파일 무결성 검사", ACCENT,
     ["SHA-256 해시 업로드 시 자동 계산 및 DB 저장",
      "버전 관리: 동일 파일명 재업로드 시 v2, v3 자동 부여",
      "중복 감지: 동일 해시 → 경고 반환",
      "무결성 검증: 실제 해시 vs 저장 해시 비교 (변조 감지)",
      "업로드 → Google Drive 자동 push → Discord 알림"]),
    ("📅  캘린더", ACCENT2,
     ["월별 그리드 뷰 (Tailwind CSS Grid, 외부 라이브러리 없음)",
      "태스크 (우선순위별 색: URGENT=자홍·HIGH=주황·MEDIUM=황·LOW=회)",
      "회의 일정 (하늘색) 동시 표시",
      "일정 추가 모달: 태스크 탭 / 회의 탭",
      "생성 즉시 Discord 알림 전송"]),
    ("🚨  스마트 경보 시스템", RGBColor(0xFB,0x7B,0x5E),
     ["IMBALANCE: 기여도 편차 과도 (Gini 계수 기반)",
      "INACTIVITY: 특정 멤버 7일 이상 활동 없음",
      "OVERLOAD: 특정 멤버 과부하 (평균 2배 이상)",
      "CRAMMING: 마감 직전 집중 활동 감지",
      "DEADLINE: 마감 D-3 이내 태스크 자동 감지"]),
]

for i, (title, color, bullets) in enumerate(feats):
    lx = 0.4 + i * 4.25
    rect(sl, lx, 1.8, 4.0, 5.35, fill=CARD_BG)
    sh = sl.shapes.add_shape(1, Inches(lx), Inches(1.8), Inches(4.0), Inches(0.07))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    txbox(sl, title, lx+0.15, 1.9, 3.7, 0.5, size=13, bold=True, color=DARK_TEXT)
    for j, b in enumerate(bullets):
        txbox(sl, f"✓  {b}", lx+0.15, 2.5+j*0.88, 3.7, 0.78, size=11, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════
# 슬라이드 13 — AI 활용 기능
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "03  프로젝트 구현", "AI 활용 기능")

# 연동 구조
rect(sl, 0.4, 1.8, 12.5, 2.6, fill=CARD_BG)
txbox(sl, "AI 연동 구조", 0.6, 1.85, 5, 0.35, size=12, bold=True, color=ACCENT)

arch_items = [
    ("Frontend\n(Next.js)", 0.6, 2.4),
    ("Backend\n(Spring Boot)", 3.2, 2.4),
    ("OpenAI\nGPT-4.1-nano", 6.2, 2.4),
    ("ai_analysis\n_results (DB)", 9.3, 2.4),
]
for label, lx, ty in arch_items:
    rect(sl, lx, ty, 2.5, 1.5, fill=RGBColor(0xFF,0xFF,0xFF), line=ACCENT2)
    txbox(sl, label, lx+0.1, ty+0.45, 2.3, 0.7, size=12, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
for i in range(3):
    txbox(sl, "→", 3.0+i*3.1, 2.9, 0.5, 0.5, size=18, color=ACCENT2)

# AI 기능 설명
ai_features = [
    ("🤖  팀 역량 분석",
     "팀 전체 활동 데이터를 GPT-4.1-nano에 전달\n→ 팀 강점/약점/개선 방향 자연어 보고서 생성\n→ DB 캐싱 (중복 API 호출 방지)"),
    ("📊  커밋 품질 분석",
     "최근 커밋 메시지·변경량 GPT 분석\n→ 1~5점 품질 점수 산출\n→ scoreWeight에 곱해 기여도 점수 반영 (x0.5~x1.5)"),
    ("🛡️  개발 보조",
     "Claude Code (claude-sonnet-4-6) 활용\n→ 전체 백엔드/프론트엔드 설계 및 구현 지원\n→ 버그 수정, 아키텍처 결정, DB 스키마 설계"),
]
for i, (title, desc) in enumerate(ai_features):
    lx = 0.4 + i * 4.25
    rect(sl, lx, 4.6, 4.0, 2.55, fill=RGBColor(0xE8,0xEF,0xF5), line=ACCENT)
    txbox(sl, title, lx+0.15, 4.7, 3.7, 0.45, size=13, bold=True, color=ACCENT)
    txbox(sl, desc,  lx+0.15, 5.2, 3.7, 1.8,  size=11, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════
# 슬라이드 14 — 2학기 진행 계획
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "04  2학기 진행 계획", "신구EXPO 발표 대비 기능 확장")

plans = [
    ("7~8월", "기반 고도화",
     ["라이브러리 시스템 구현 (참고자료 공유 기능)",
      "백업 플로우 UI 추가",
      "교수 대시보드 강화 (전체 프로젝트 조망)"],
     ACCENT),
    ("9월", "사용성 개선",
     ["반응형 모바일 UI 최적화",
      "알림 설정 세분화 (채널별 온/오프)",
      "성능 최적화 (API 캐싱, DB 인덱싱)"],
     ACCENT2),
    ("10월", "테스트 & 안정화",
     ["실제 수업 적용 베타 테스트",
      "사용자 피드백 반영",
      "신구EXPO 발표 준비 및 데모 환경 구축"],
     RGBColor(0xFB,0x7B,0x5E)),
]
for i, (period, phase, tasks, color) in enumerate(plans):
    lx = 0.4 + i * 4.25
    rect(sl, lx, 1.8, 4.0, 4.8, fill=CARD_BG)
    ph_sh = sl.shapes.add_shape(1, Inches(lx), Inches(1.8), Inches(4.0), Inches(0.08))
    ph_sh.fill.solid(); ph_sh.fill.fore_color.rgb = color; ph_sh.line.fill.background()
    txbox(sl, period, lx+0.15, 1.92, 3.7, 0.38, size=14, bold=True, color=color)
    txbox(sl, phase,  lx+0.15, 2.35, 3.7, 0.38, size=12, color=DARK_TEXT)
    for j, t in enumerate(tasks):
        txbox(sl, f"• {t}", lx+0.15, 2.85+j*0.95, 3.7, 0.85, size=11, color=LIGHT_GRAY)

rect(sl, 0.4, 6.75, 12.5, 0.55, fill=RGBColor(0xEC,0xFD,0xF5))
txbox(sl, "목표: 신구EXPO 전시 가능한 완성도 높은 서비스로 고도화",
      0.6, 6.82, 12.0, 0.38, size=12, bold=True, color=ACCENT)

# ══════════════════════════════════════════════════════════
# 슬라이드 15 — 결과 및 고찰
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "05  결과 및 고찰")

# 소감 박스
rect(sl, 0.4, 1.78, 12.5, 1.9, fill=CARD_BG)
txbox(sl, "프로젝트 소감 및 느낀 점", 0.6, 1.83, 10, 0.35, size=13, bold=True, color=ACCENT)
txbox(sl,
      "• 외부 API(GitHub·Notion·Google·Discord·OpenAI) 연동을 직접 구현하며 OAuth 2.0 인증 흐름과 Webhook 설계를 실전으로 익혔음\n"
      "• Spring Boot + Next.js 풀스택과 Docker Compose 배포 환경을 처음부터 설계·구축하며 DevOps 기초를 체험\n"
      "• 기여도를 단순 수치가 아닌 AI 분석으로 정성적 진단까지 연결하는 시스템 설계의 중요성을 깨달음",
      0.6, 2.22, 12.1, 1.35, size=11.5, color=DARK_TEXT)

# 시행착오 제목
txbox(sl, "시행착오 및 극복 사례", 0.4, 3.78, 6, 0.35, size=13, bold=True, color=RGBColor(0xFB,0x7B,0x5E))

challenges = [
    ("① GitHub OAuth 연동 실패",
     "문제: GitHub 로그인 클릭 시\nclient_id 빈값 오류 발생\n\n원인: .env 값이 docker-compose.yml\n환경변수 블록에 누락됨\n\n해결: docker-compose.yml에\nGITHUB_CLIENT_ID 명시적 추가"),
    ("② Java Map.of() 타입 오류",
     "문제: Map.of(\"sha\", str, \"score\", int)\n혼합 타입으로 컴파일 에러 발생\n\n원인: 제네릭 타입 추론 실패\n(Map<String,Object> 불가)\n\n해결: HashMap으로 교체 후\n명시적 put() 사용"),
    ("③ PDF 한글 깨짐",
     "문제: PDF 보고서 출력 시\n한글이 모두 깨져서 표시됨\n\n원인: iText 기본 폰트가\n한글 글리프 미포함\n\n해결: NanumGothic.ttf를\nDocker 이미지에 포함 후\nBaseFont로 명시 지정"),
]
for i, (title, desc) in enumerate(challenges):
    lx = 0.4 + i * 4.25
    rect(sl, lx, 4.18, 4.0, 3.05, fill=RGBColor(0xFF,0xF7,0xF5), line=RGBColor(0xFB,0x7B,0x5E))
    # 상단 컬러 바
    top = sl.shapes.add_shape(1, Inches(lx), Inches(4.18), Inches(4.0), Inches(0.06))
    top.fill.solid(); top.fill.fore_color.rgb = RGBColor(0xFB,0x7B,0x5E); top.line.fill.background()
    txbox(sl, title, lx+0.15, 4.27, 3.7, 0.42, size=12, bold=True, color=RGBColor(0xFB,0x7B,0x5E))
    txbox(sl, desc,  lx+0.15, 4.72, 3.7, 2.4,  size=10.5, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════
# 슬라이드 — 참고자료
# ══════════════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "참고자료")

refs = [
    ("📚 서적",
     ["스프링 부트 핵심 가이드 (위키북스, 장정우)",
      "자바의 정석 3rd Edition (남궁 성)"]),
    ("🎥 동영상·강의",
     ["인프런 - 스프링 입문 (김영한)",
      "Next.js 공식 문서 튜토리얼 (nextjs.org/learn)"]),
    ("🌐 공식 문서 & API",
     ["Spring Boot Reference Docs — docs.spring.io",
      "GitHub REST API Docs — docs.github.com/en/rest",
      "Notion API Docs — developers.notion.com",
      "Google Drive API Docs — developers.google.com/drive",
      "Discord Webhook Guide — discord.com/developers/docs",
      "OpenAI API Reference — platform.openai.com/docs"]),
    ("🤖 AI 활용",
     ["Claude Code (claude-sonnet-4-6) — 전체 백엔드·프론트엔드 설계 및 구현 지원,\n     버그 수정, DB 스키마 설계, Docker 환경 구성",
      "ChatGPT (GPT-4o) — 알고리즘 아이디어 검토, 오류 원인 분석"]),
]

ty = 1.85
for category, items in refs:
    rect(sl, 0.4, ty, 12.5, 0.3, fill=RGBColor(0xE8,0xEF,0xF5))
    txbox(sl, category, 0.6, ty+0.04, 12, 0.24, size=12, bold=True, color=ACCENT)
    ty += 0.34
    for item in items:
        txbox(sl, f"  • {item}", 0.6, ty, 12.1, 0.38, size=11, color=LIGHT_GRAY)
        ty += 0.38
    ty += 0.12

# ══════════════════════════════════════════════════════════
# 슬라이드 — 마무리
# ══════════════════════════════════════════════════════════
sl = add_slide()

for (cx, cy, sz, clr) in [(9, 3, 5, ACCENT), (2, 5, 3, ACCENT2)]:
    sh = sl.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = clr; sh.line.fill.background()

txbox(sl, "감사합니다", 1.0, 2.0, 11, 1.5, size=52, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
txbox(sl, "Team Blackbox  —  컴퓨터소프트웨어과  프로젝트 구현",
      1.0, 3.7, 11, 0.6, size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

highlights = ["✓  풀스택 EdTech SaaS (Spring Boot + Next.js)",
              "✓  5개 외부 서비스 연동 (GitHub·Notion·Google·Discord·OpenAI)",
              "✓  AI 기반 자동 기여도 측정 & 경보 시스템"]
for i, h in enumerate(highlights):
    txbox(sl, h, 2.5, 4.6+i*0.55, 8.5, 0.5, size=13, color=ACCENT, align=PP_ALIGN.CENTER)

# ── 저장 ─────────────────────────────────────────────────
out = r"C:\blackbox\Blackbox_presentation.pptx"
prs.save(out)
print(f"저장 완료: {out}")
print(f"슬라이드 수: {len(prs.slides)}")
